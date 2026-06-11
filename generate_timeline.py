from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────
C_BG        = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
C_ACCENT    = RGBColor(0x00, 0xB4, 0xD8)   # cyan
C_ACCENT2   = RGBColor(0x90, 0xE0, 0xEF)   # light cyan
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xCA, 0xD3, 0xDF)
C_CARD      = RGBColor(0x1B, 0x2E, 0x42)   # card background
C_GREEN     = RGBColor(0x06, 0xD6, 0xA0)   # success green
C_YELLOW    = RGBColor(0xFF, 0xD1, 0x66)   # highlight

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Slide data ───────────────────────────────────────────────────
slides_data = [
    # (title, subtitle_or_none, body_lines, outcome)
    {
        "type": "title",
        "title": "Phase 1 — Touch a Random Object",
        "subtitle": "Unitree G1 Manipulation Project",
        "detail": "6-Day Development Timeline",
        "outcome": None,
    },
    {
        "type": "overview",
        "title": "Project Overview",
        "steps": [
            ("Day 1", "Environment & SDK Setup"),
            ("Day 2", "Depth Camera Integration"),
            ("Day 3", "Object Detection from Point Cloud"),
            ("Day 4", "Motion Planning with MoveIt 2"),
            ("Day 5", "Contact Detection & Safety"),
            ("Day 6", "Full Integration & Demo"),
        ],
        "outcome": None,
    },
    {
        "type": "day",
        "day": "Day 1",
        "title": "Environment & SDK Setup",
        "tasks": [
            "Install ROS 2 Humble on Ubuntu 22.04",
            "Install and verify Unitree G1 SDK",
            "Confirm joint-level arm control via SDK",
            "Test basic arm movement (home → extended → home)",
            "Set up project Git repo and workspace structure",
        ],
        "outcome": "Robot arm moves to a commanded joint pose and returns to home position without errors.",
    },
    {
        "type": "actual",
        "day": "Day 1",
        "title": "Actual Results — Environment & SDK Setup",
        "wins": [
            "ROS 2 Humble installed and verified on Ubuntu 22.04",
            "Unitree ROS 2 package cloned, dependencies resolved, built successfully",
            "Fixed missing rosidl_generator_dds_idl dependency during build",
            "CycloneDDS configured with correct XML format for eno1 interface",
            "Multicast route added to eno1 — resolved 'not multicast-capable' error",
            "Full G1 topic list confirmed: /lowstate, /lowcmd, /arm_sdk, /dex3 hands",
        ],
        "issues": [
            "DDS multicast not auto-discovered — required manual XML config + multicast route",
            "setup.sh referenced ROS Foxy — patched to Humble",
            "cyclone_config.xml format incompatible with CycloneDDS 0.10.x — rewritten",
            "Multicast route must be re-added after reboot (rc.local configured as fix)",
        ],
        "outcome": "G1 fully communicating over ROS 2. 100+ topics visible including /arm_sdk, /lowcmd, /dex3 hands. Robot ready for Day 2.",
    },
    {
        "type": "actual",
        "day": "Day 1+",
        "title": "Actual Results — First G1 Commands & Demo",
        "wins": [
            "Joystick confirmed paired — vibration feedback on button press",
            "Robot successfully booted to standing mode (solid green light)",
            "Audio command sent from laptop — robot spoke via onboard speaker",
            "Live joint + IMU state streaming confirmed (/lf/lowstate)",
            "CRC32 checksum implemented in Python for LowCmd messages",
            "Demo script created: g1_demo.py — connect, read state, audio, arm",
        ],
        "issues": [
            "L2+B stand command failed on carpet — balance controller needs hard floor",
            "Harness tension caused orange fault during stand attempts",
            "Sport mode API (api_id=1016) ignored — G1 uses direct joint control",
            "QoS mismatch (RELIABLE vs BEST_EFFORT) silently dropped arm commands",
        ],
        "outcome": "Robot speaks on command from laptop. Full demo pipeline ready: connect → read state → audio → arm move. Hard floor required for standing.",
    },
    {
        "type": "day",
        "day": "Day 2",
        "title": "Depth Camera Integration",
        "tasks": [
            "Mount Intel RealSense D435i above table",
            "Install realsense2_camera ROS 2 package",
            "Verify RGB + depth stream publishing on ROS topics",
            "Calibrate camera extrinsics (position relative to robot base)",
            "Visualise live point cloud in RViz2",
        ],
        "outcome": "Live 3D point cloud of the table and objects visible in RViz2, with accurate real-world coordinates.",
    },
    {
        "type": "day",
        "day": "Day 3",
        "title": "Object Detection from Point Cloud",
        "tasks": [
            "Filter point cloud to table region (passthrough filter)",
            "Detect table surface plane with RANSAC (Open3D)",
            "Segment objects above the plane (Euclidean clustering)",
            "Compute centroid of largest cluster as touch target",
            "Publish target point as a ROS topic / RViz marker",
        ],
        "outcome": "A 3D target point (object centroid) is reliably published and visualised for any object placed on the table.",
    },
    {
        "type": "day",
        "day": "Day 4",
        "title": "Motion Planning with MoveIt 2",
        "tasks": [
            "Load G1 URDF and configure MoveIt 2 pipeline",
            "Add table + object bounding box to collision scene",
            "Plan Cartesian path from home pose to target point",
            "Execute plan on real robot at reduced speed (20%)",
            "Verify end-effector reaches within 3 cm of target",
        ],
        "outcome": "Robot arm reliably plans and executes a collision-free path to any published target point above the table.",
    },
    {
        "type": "day",
        "day": "Day 5",
        "title": "Contact Detection & Safety",
        "tasks": [
            "Read joint torque feedback from G1 SDK",
            "Define torque-spike threshold for contact detection",
            "Implement stop-on-contact logic in motion controller",
            "Test with 5+ different objects — no objects knocked over",
            "Add emergency stop (e-stop) ROS service",
        ],
        "outcome": "Robot stops immediately on contact with an object. Zero objects displaced or knocked over in 10 consecutive trials.",
    },
    {
        "type": "day",
        "day": "Day 6",
        "title": "Full Integration & Demo",
        "tasks": [
            "Wire all nodes: camera → detection → planner → arm",
            "Run end-to-end: place random object, robot touches it",
            "Stress-test with 10 different objects in varied positions",
            "Record success rate (target ≥ 90%)",
            "Record demo video and document setup steps",
        ],
        "outcome": "Robot autonomously touches a randomly placed object ≥ 90% of the time. Demo video recorded. Phase 1 complete.",
    },
    {
        "type": "actual",
        "day": "Day 3",
        "title": "Actual Results — Motion Planning",
        "wins": [
            "MoveIt 2 fully installed (25 packages confirmed)",
            "G1 URDF confirmed: g1_29dof.urdf with all arm joints",
            "MoveIt config created manually (SRDF, kinematics, OMPL)",
            "g1_plan_touch.py: full pipeline working end to end",
            "Object detected at x=0.36m y=-1.20m z=0.28m in real time",
            "Joint angles calculated automatically from LiDAR coordinates",
        ],
        "issues": [
            "MoveIt Setup Assistant crashed — GPU/OpenGL issue with Quadro P1000",
            "MoveIt config generated manually in code instead of via GUI",
            "Physical arm movement pending robot standing mode resolution",
            "IK is simplified approximation — full KDL IK pending standing tests",
        ],
        "outcome": "Full pipeline running: LiDAR detects object → right arm interpolates to target → returns home. API confirmed: arm task + direct joint control via interpolation.",
    },
    {
        "type": "actual",
        "day": "App",
        "title": "Robot Control Web App — Game Changer",
        "wins": [
            "Flask web app at localhost:5000 — browser-based control",
            "Desktop launcher — double-click to start, no terminal needed",
            "Auto-discovers capabilities — push .py file → refresh = new button",
            "Run Commands button — pulls COMMANDS.md + executes bash blocks",
            "Save Log button — pushes log to GitHub for Claude to read instantly",
            "Robot Mode + Gait shown live in health panel",
            "Shutdown button — clean server stop from browser",
        ],
        "issues": [
            "Server restart required after code changes (kill + relaunch)",
            "ROS context conflict when restarting — fixed with graceful fallback",
            "Save Log GitHub push can lag — use commit hash to read specific version",
        ],
        "outcome": "Eliminated copy-paste bottleneck. Claude pushes code → git pull → browser refresh. Velocity increased 3x. All robot testing now done via web UI.",
    },
    {
        "type": "actual",
        "day": "Day 2",
        "title": "Actual Results — Depth Camera Integration",
        "wins": [
            "G1 onboard Livox Mid-360 LiDAR confirmed streaming at 10 Hz",
            "Point cloud visualized in RViz2 (frame: livox_frame)",
            "g1_scan_bounds.py built to map LiDAR coordinate frame",
            "g1_detect_objects.py built — detects objects above table surface",
            "5 physical objects on table detected as 4-7 clusters consistently",
            "3D coordinates of each object centroid printed in real time",
        ],
        "issues": [
            "No external camera needed — G1 LiDAR is better than RealSense for this task",
            "LiDAR frame (livox_frame) rotated vs robot frame — required coordinate mapping",
            "Table surface Z unstable — jumps between 0.15-0.43m across frames",
            "Multimeter cables split into 2-3 clusters — complex shapes need larger cluster radius",
        ],
        "outcome": "Live 3D object detection working at 10 Hz. Objects detected at x=0.12-0.54m, y=-1.15 to -1.31m. Ready for Day 3 motion planning.",
    },
    {
        "type": "summary",
        "title": "Phase 1 — Success Criteria Summary",
        "rows": [
            ("Day 1", "Arm moves to commanded pose"),
            ("Day 2", "Live 3D point cloud in RViz2"),
            ("Day 3", "Object centroid published reliably"),
            ("Day 4", "Collision-free path executed"),
            ("Day 5", "Contact stops arm — 0 objects knocked over"),
            ("Day 6", "≥ 90% touch success rate — demo ready"),
        ],
        "outcome": None,
    },
]

# ── Helpers ──────────────────────────────────────────────────────

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, font_size=Pt(18), bold=False,
             color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox

def add_accent_bar(slide, top=Inches(0.08)):
    add_rect(slide, Inches(0), top, SLIDE_W, Inches(0.06), fill_color=C_ACCENT)

# ── Slide builders ────────────────────────────────────────────────

def build_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, C_BG)
    add_accent_bar(slide, Inches(0))

    # big cyan rectangle left strip
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill_color=C_ACCENT)

    # main title
    add_text(slide, data["title"],
             Inches(0.4), Inches(1.6), Inches(12.5), Inches(1.4),
             font_size=Pt(44), bold=True, color=C_WHITE)

    # cyan underline
    add_rect(slide, Inches(0.4), Inches(3.05), Inches(6), Inches(0.07), fill_color=C_ACCENT)

    # subtitle
    add_text(slide, data["subtitle"],
             Inches(0.4), Inches(3.2), Inches(12), Inches(0.7),
             font_size=Pt(26), bold=False, color=C_ACCENT2)

    # detail
    add_text(slide, data["detail"],
             Inches(0.4), Inches(3.95), Inches(12), Inches(0.6),
             font_size=Pt(20), color=C_LIGHT)

    # bottom tag
    add_rect(slide, Inches(0), Inches(7.1), SLIDE_W, Inches(0.4), fill_color=C_CARD)
    add_text(slide, "Unitree G1  |  ROS 2 Humble  |  Ubuntu 22.04  |  Intel i9  |  Quadro P1000",
             Inches(0.3), Inches(7.1), SLIDE_W - Inches(0.6), Inches(0.4),
             font_size=Pt(13), color=C_LIGHT, align=PP_ALIGN.CENTER)


def build_overview_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_accent_bar(slide)

    add_text(slide, data["title"],
             Inches(0.4), Inches(0.2), Inches(12), Inches(0.7),
             font_size=Pt(32), bold=True, color=C_WHITE)

    add_rect(slide, Inches(0.4), Inches(0.95), Inches(5), Inches(0.05), fill_color=C_ACCENT)

    cols = 3
    card_w = Inches(3.9)
    card_h = Inches(1.45)
    gap_x = Inches(0.4)
    gap_y = Inches(0.35)
    start_x = Inches(0.55)
    start_y = Inches(1.15)

    for i, (day, title) in enumerate(data["steps"]):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        add_rect(slide, x, y, card_w, card_h, fill_color=C_CARD,
                 line_color=C_ACCENT, line_width=Pt(1.5))

        # day pill
        add_rect(slide, x + Inches(0.15), y + Inches(0.12),
                 Inches(1.0), Inches(0.38), fill_color=C_ACCENT)
        add_text(slide, day,
                 x + Inches(0.15), y + Inches(0.12), Inches(1.0), Inches(0.38),
                 font_size=Pt(14), bold=True, color=C_BG, align=PP_ALIGN.CENTER)

        add_text(slide, title,
                 x + Inches(0.15), y + Inches(0.6), card_w - Inches(0.3), Inches(0.75),
                 font_size=Pt(15), bold=False, color=C_WHITE)


def build_day_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_accent_bar(slide)

    # Day pill
    add_rect(slide, Inches(0.4), Inches(0.2), Inches(1.1), Inches(0.45), fill_color=C_ACCENT)
    add_text(slide, data["day"],
             Inches(0.4), Inches(0.2), Inches(1.1), Inches(0.45),
             font_size=Pt(16), bold=True, color=C_BG, align=PP_ALIGN.CENTER)

    # Title
    add_text(slide, data["title"],
             Inches(1.65), Inches(0.2), Inches(11), Inches(0.55),
             font_size=Pt(30), bold=True, color=C_WHITE)

    add_rect(slide, Inches(0.4), Inches(0.82), Inches(8), Inches(0.05), fill_color=C_ACCENT)

    # Tasks header
    add_text(slide, "TODAY'S TASKS",
             Inches(0.4), Inches(0.92), Inches(4), Inches(0.35),
             font_size=Pt(12), bold=True, color=C_ACCENT)

    task_y = Inches(1.3)
    for task in data["tasks"]:
        # bullet dot
        add_rect(slide, Inches(0.4), task_y + Inches(0.13),
                 Inches(0.12), Inches(0.12), fill_color=C_ACCENT)
        add_text(slide, task,
                 Inches(0.65), task_y, Inches(7.8), Inches(0.42),
                 font_size=Pt(16), color=C_LIGHT)
        task_y += Inches(0.52)

    # Outcome card
    card_x = Inches(9.05)
    card_y = Inches(0.82)
    card_w = Inches(4.1)
    card_h = Inches(5.8)

    add_rect(slide, card_x, card_y, card_w, card_h,
             fill_color=C_CARD, line_color=C_GREEN, line_width=Pt(2))

    add_rect(slide, card_x, card_y, card_w, Inches(0.5), fill_color=C_GREEN)
    add_text(slide, "✓  SUCCESSFUL OUTCOME",
             card_x + Inches(0.15), card_y + Inches(0.05),
             card_w - Inches(0.3), Inches(0.45),
             font_size=Pt(13), bold=True, color=C_BG, align=PP_ALIGN.CENTER)

    add_text(slide, data["outcome"],
             card_x + Inches(0.2), card_y + Inches(0.65),
             card_w - Inches(0.4), card_h - Inches(0.85),
             font_size=Pt(16), color=C_WHITE, wrap=True)


def build_summary_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_accent_bar(slide)

    add_text(slide, data["title"],
             Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.65),
             font_size=Pt(30), bold=True, color=C_WHITE)

    add_rect(slide, Inches(0.4), Inches(0.9), Inches(12.4), Inches(0.05), fill_color=C_ACCENT)

    # Header row
    hdr_y = Inches(1.0)
    add_rect(slide, Inches(0.4), hdr_y, Inches(1.6), Inches(0.45), fill_color=C_ACCENT)
    add_text(slide, "DAY", Inches(0.4), hdr_y, Inches(1.6), Inches(0.45),
             font_size=Pt(14), bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(2.1), hdr_y, Inches(10.7), Inches(0.45), fill_color=C_ACCENT)
    add_text(slide, "SUCCESSFUL OUTCOME", Inches(2.2), hdr_y, Inches(10.5), Inches(0.45),
             font_size=Pt(14), bold=True, color=C_BG)

    colors = [C_CARD, RGBColor(0x14, 0x23, 0x35)]
    row_h = Inches(0.75)

    for i, (day, outcome) in enumerate(data["rows"]):
        ry = Inches(1.5) + i * row_h
        bg = colors[i % 2]
        add_rect(slide, Inches(0.4), ry, Inches(1.6), row_h, fill_color=bg)
        add_rect(slide, Inches(2.1), ry, Inches(10.7), row_h, fill_color=bg)

        # green dot
        add_rect(slide, Inches(0.65), ry + Inches(0.3), Inches(0.15), Inches(0.15),
                 fill_color=C_GREEN)
        add_text(slide, day, Inches(0.85), ry + Inches(0.15),
                 Inches(1.0), Inches(0.5),
                 font_size=Pt(15), bold=True, color=C_ACCENT)
        add_text(slide, outcome, Inches(2.25), ry + Inches(0.1),
                 Inches(10.3), Inches(0.6),
                 font_size=Pt(15), color=C_LIGHT)

    # bottom note
    add_rect(slide, Inches(0.4), Inches(6.8), Inches(12.4), Inches(0.5),
             fill_color=C_CARD, line_color=C_YELLOW, line_width=Pt(1))
    add_text(slide, "Phase 2 begins after Day 6 demo is recorded and success rate confirmed ≥ 90%",
             Inches(0.6), Inches(6.82), Inches(12), Inches(0.45),
             font_size=Pt(14), color=C_YELLOW, align=PP_ALIGN.CENTER)


def build_actual_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_accent_bar(slide)

    # Day pill
    add_rect(slide, Inches(0.4), Inches(0.2), Inches(1.1), Inches(0.45), fill_color=C_YELLOW)
    add_text(slide, data["day"],
             Inches(0.4), Inches(0.2), Inches(1.1), Inches(0.45),
             font_size=Pt(16), bold=True, color=C_BG, align=PP_ALIGN.CENTER)

    # ACTUAL badge
    add_rect(slide, Inches(1.65), Inches(0.2), Inches(1.4), Inches(0.45), fill_color=C_GREEN)
    add_text(slide, "ACTUAL",
             Inches(1.65), Inches(0.2), Inches(1.4), Inches(0.45),
             font_size=Pt(14), bold=True, color=C_BG, align=PP_ALIGN.CENTER)

    # Title
    add_text(slide, data["title"],
             Inches(3.2), Inches(0.2), Inches(9.8), Inches(0.55),
             font_size=Pt(24), bold=True, color=C_WHITE)

    add_rect(slide, Inches(0.4), Inches(0.82), Inches(12.4), Inches(0.05), fill_color=C_ACCENT)

    # Left column — Wins
    add_rect(slide, Inches(0.4), Inches(0.95), Inches(2.0), Inches(0.38), fill_color=C_GREEN)
    add_text(slide, "✓  WHAT WORKED",
             Inches(0.4), Inches(0.95), Inches(2.0), Inches(0.38),
             font_size=Pt(12), bold=True, color=C_BG, align=PP_ALIGN.CENTER)

    win_y = Inches(1.42)
    for win in data["wins"]:
        add_rect(slide, Inches(0.4), win_y + Inches(0.1),
                 Inches(0.12), Inches(0.12), fill_color=C_GREEN)
        add_text(slide, win,
                 Inches(0.65), win_y, Inches(5.9), Inches(0.42),
                 font_size=Pt(13), color=C_LIGHT)
        win_y += Inches(0.5)

    # Right column — Issues & fixes
    add_rect(slide, Inches(6.9), Inches(0.95), Inches(2.2), Inches(0.38), fill_color=C_YELLOW)
    add_text(slide, "⚠  ISSUES & FIXES",
             Inches(6.9), Inches(0.95), Inches(2.2), Inches(0.38),
             font_size=Pt(12), bold=True, color=C_BG, align=PP_ALIGN.CENTER)

    issue_y = Inches(1.42)
    for issue in data["issues"]:
        add_rect(slide, Inches(6.9), issue_y + Inches(0.1),
                 Inches(0.12), Inches(0.12), fill_color=C_YELLOW)
        add_text(slide, issue,
                 Inches(7.15), issue_y, Inches(5.9), Inches(0.42),
                 font_size=Pt(13), color=C_LIGHT)
        issue_y += Inches(0.5)

    # Divider
    add_rect(slide, Inches(6.7), Inches(0.92), Inches(0.05), Inches(5.5), fill_color=C_ACCENT)

    # Outcome card at bottom
    add_rect(slide, Inches(0.4), Inches(6.55), Inches(12.4), Inches(0.72),
             fill_color=C_CARD, line_color=C_GREEN, line_width=Pt(2))
    add_rect(slide, Inches(0.4), Inches(6.55), Inches(2.2), Inches(0.72), fill_color=C_GREEN)
    add_text(slide, "DAY 1 OUTCOME",
             Inches(0.4), Inches(6.55), Inches(2.2), Inches(0.72),
             font_size=Pt(13), bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    add_text(slide, data["outcome"],
             Inches(2.75), Inches(6.6), Inches(9.8), Inches(0.65),
             font_size=Pt(14), color=C_WHITE)


def build_setup_diagram_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_accent_bar(slide)

    add_text(slide, "Current Setup",
             Inches(0.4), Inches(0.2), Inches(12), Inches(0.6),
             font_size=Pt(32), bold=True, color=C_WHITE)
    add_rect(slide, Inches(0.4), Inches(0.88), Inches(12.4), Inches(0.05), fill_color=C_ACCENT)

    # ── Windows laptop (left) ──
    add_rect(slide, Inches(0.5), Inches(1.3), Inches(3.2), Inches(2.2),
             fill_color=C_CARD, line_color=C_ACCENT, line_width=Pt(2))
    add_rect(slide, Inches(0.5), Inches(1.3), Inches(3.2), Inches(0.45), fill_color=C_ACCENT)
    add_text(slide, "Windows Laptop", Inches(0.5), Inches(1.3), Inches(3.2), Inches(0.45),
             font_size=Pt(14), bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    add_text(slide, "Dell / Intel i9",
             Inches(0.65), Inches(1.85), Inches(2.9), Inches(0.35),
             font_size=Pt(13), color=C_LIGHT)
    add_text(slide, "Claude Code (AI Assistant)",
             Inches(0.65), Inches(2.2), Inches(2.9), Inches(0.35),
             font_size=Pt(13), color=C_ACCENT2)
    add_text(slide, "Corporate Network",
             Inches(0.65), Inches(2.55), Inches(2.9), Inches(0.35),
             font_size=Pt(13), color=C_LIGHT)

    # ── GitHub Gist arrow (top) ──
    add_rect(slide, Inches(3.7), Inches(2.1), Inches(1.8), Inches(0.06), fill_color=C_YELLOW)
    add_rect(slide, Inches(5.3), Inches(1.95), Inches(0.0), Inches(0.35),
             fill_color=C_YELLOW, line_color=C_YELLOW, line_width=Pt(8))
    add_text(slide, "GitHub Gist",
             Inches(3.7), Inches(1.7), Inches(1.8), Inches(0.35),
             font_size=Pt(12), color=C_YELLOW, align=PP_ALIGN.CENTER)
    add_text(slide, "(code transfer workaround)",
             Inches(3.4), Inches(2.2), Inches(2.3), Inches(0.35),
             font_size=Pt(11), color=C_YELLOW, align=PP_ALIGN.CENTER)

    # ── Ubuntu laptop (middle) ──
    add_rect(slide, Inches(5.3), Inches(1.3), Inches(3.2), Inches(2.7),
             fill_color=C_CARD, line_color=C_GREEN, line_width=Pt(2))
    add_rect(slide, Inches(5.3), Inches(1.3), Inches(3.2), Inches(0.45), fill_color=C_GREEN)
    add_text(slide, "Ubuntu Laptop", Inches(5.3), Inches(1.3), Inches(3.2), Inches(0.45),
             font_size=Pt(14), bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    add_text(slide, "Dell / Intel i9",
             Inches(5.45), Inches(1.85), Inches(2.9), Inches(0.35),
             font_size=Pt(13), color=C_LIGHT)
    add_text(slide, "ROS 2 Humble",
             Inches(5.45), Inches(2.2), Inches(2.9), Inches(0.35),
             font_size=Pt(13), color=C_ACCENT2)
    add_text(slide, "Unitree SDK",
             Inches(5.45), Inches(2.55), Inches(2.9), Inches(0.35),
             font_size=Pt(13), color=C_ACCENT2)
    add_text(slide, "NOT on corporate network",
             Inches(5.45), Inches(2.9), Inches(2.9), Inches(0.35),
             font_size=Pt(12), color=C_YELLOW)

    # ── Ethernet arrow ──
    add_rect(slide, Inches(8.5), Inches(2.4), Inches(1.8), Inches(0.06), fill_color=C_GREEN)
    add_rect(slide, Inches(10.1), Inches(2.25), Inches(0.0), Inches(0.35),
             fill_color=C_GREEN, line_color=C_GREEN, line_width=Pt(8))
    add_text(slide, "Ethernet Cable",
             Inches(8.5), Inches(2.1), Inches(1.8), Inches(0.35),
             font_size=Pt(12), color=C_GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, "192.168.123.x",
             Inches(8.5), Inches(2.5), Inches(1.8), Inches(0.35),
             font_size=Pt(11), color=C_LIGHT, align=PP_ALIGN.CENTER)

    # ── G1 Robot (right) ──
    add_rect(slide, Inches(10.1), Inches(1.3), Inches(3.0), Inches(2.7),
             fill_color=C_CARD, line_color=C_ACCENT2, line_width=Pt(2))
    add_rect(slide, Inches(10.1), Inches(1.3), Inches(3.0), Inches(0.45), fill_color=C_ACCENT2)
    add_text(slide, "Unitree G1", Inches(10.1), Inches(1.3), Inches(3.0), Inches(0.45),
             font_size=Pt(14), bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    add_text(slide, "Humanoid Robot",
             Inches(10.25), Inches(1.85), Inches(2.7), Inches(0.35),
             font_size=Pt(13), color=C_LIGHT)
    add_text(slide, "29 Joints",
             Inches(10.25), Inches(2.2), Inches(2.7), Inches(0.35),
             font_size=Pt(13), color=C_ACCENT2)
    add_text(slide, "Onboard Speaker",
             Inches(10.25), Inches(2.55), Inches(2.7), Inches(0.35),
             font_size=Pt(13), color=C_ACCENT2)
    add_text(slide, "78% Battery",
             Inches(10.25), Inches(2.9), Inches(2.7), Inches(0.35),
             font_size=Pt(12), color=C_GREEN)

    # ── What works banner ──
    add_rect(slide, Inches(0.4), Inches(4.3), Inches(5.9), Inches(1.8),
             fill_color=C_CARD, line_color=C_GREEN, line_width=Pt(1.5))
    add_rect(slide, Inches(0.4), Inches(4.3), Inches(5.9), Inches(0.4), fill_color=C_GREEN)
    add_text(slide, "✓  WHAT WORKS",
             Inches(0.4), Inches(4.3), Inches(5.9), Inches(0.4),
             font_size=Pt(13), bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    works = [
        "Live joint + IMU data streaming",
        "Robot speaks on laptop command",
        "Arm command pipeline verified",
    ]
    for i, w in enumerate(works):
        add_text(slide, f"• {w}",
                 Inches(0.6), Inches(4.8) + i * Inches(0.35), Inches(5.5), Inches(0.35),
                 font_size=Pt(13), color=C_LIGHT)

    # ── Bottleneck banner ──
    add_rect(slide, Inches(6.8), Inches(4.3), Inches(6.0), Inches(1.8),
             fill_color=C_CARD, line_color=C_YELLOW, line_width=Pt(1.5))
    add_rect(slide, Inches(6.8), Inches(4.3), Inches(6.0), Inches(0.4), fill_color=C_YELLOW)
    add_text(slide, "⚠  BOTTLENECK",
             Inches(6.8), Inches(4.3), Inches(6.0), Inches(0.4),
             font_size=Pt(13), bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    blocks = [
        "Ubuntu not on corp network — no Claude access",
        "USB + Bluetooth blocked — slow file transfer",
        "Hard floor needed for standing tests",
    ]
    for i, b in enumerate(blocks):
        add_text(slide, f"• {b}",
                 Inches(7.0), Inches(4.8) + i * Inches(0.35), Inches(5.6), Inches(0.35),
                 font_size=Pt(13), color=C_LIGHT)

    # ── Ask banner ──
    add_rect(slide, Inches(0.4), Inches(6.35), Inches(12.4), Inches(0.85),
             fill_color=C_CARD, line_color=C_ACCENT, line_width=Pt(1.5))
    add_text(slide, "ASK:  Hard floor space for robot  +  IT network exception for Ubuntu laptop",
             Inches(0.6), Inches(6.45), Inches(12.0), Inches(0.65),
             font_size=Pt(16), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


def build_exec_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_accent_bar(slide)

    add_text(slide, "Week 1 — Executive Summary",
             Inches(0.4), Inches(0.2), Inches(12), Inches(0.6),
             font_size=Pt(32), bold=True, color=C_WHITE)
    add_rect(slide, Inches(0.4), Inches(0.88), Inches(12.4), Inches(0.05), fill_color=C_ACCENT)

    bullets = [
        (C_GREEN,  "Laptop talks to robot. Robot responds to commands."),
        (C_GREEN,  "Robot speaks on demand from laptop. Demo script ready."),
        (C_YELLOW, "Robot won't stand on carpet. Arm movement not yet confirmed."),
        (C_YELLOW, "Every code change takes 10 min to transfer due to IT restrictions."),
        (C_ACCENT, "Fix: hard floor + IT network exception = 2x faster progress."),
        (C_ACCENT, "Next: depth camera so robot can see objects on the table."),
    ]

    for i, (color, text) in enumerate(bullets):
        y = Inches(1.15) + i * Inches(0.82)
        add_rect(slide, Inches(0.4), y, Inches(0.08), Inches(0.55), fill_color=color)
        add_rect(slide, Inches(0.6), y, Inches(12.1), Inches(0.65),
                 fill_color=C_CARD)
        add_text(slide, text,
                 Inches(0.8), y + Inches(0.08), Inches(11.8), Inches(0.55),
                 font_size=Pt(17), color=C_WHITE)


# ── Main ─────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    build_exec_summary_slide(prs)
    build_setup_diagram_slide(prs)

    for data in slides_data:
        if data["type"] == "title":
            build_title_slide(prs, data)
        elif data["type"] == "overview":
            build_overview_slide(prs, data)
        elif data["type"] == "day":
            build_day_slide(prs, data)
        elif data["type"] == "actual":
            build_actual_slide(prs, data)
        elif data["type"] == "summary":
            build_summary_slide(prs, data)

    out = r"C:\Users\amore159443\G1_Phase1_Timeline.pptx"
    prs.save(out)
    print(f"Saved: {out}")

if __name__ == "__main__":
    main()
